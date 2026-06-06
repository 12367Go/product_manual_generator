# -*- coding: utf-8 -*-
"""
PPT模板解析与生成模块
负责解析用户上传的PPT模板，提取结构信息，并生成最终的产品手册

核心设计原则：
1. 直接使用原始模板文件创建新Presentation，保留所有原生格式
2. 修改单元格文本时，只替换文字内容，保留原有字体格式（除非用户指定了新格式）
3. 图片插入位置基于模板中款号单元格的位置计算
"""

import os
import copy
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import etree

from config import TEMPLATE_LABELS, DEFAULT_IMAGE_SIZE


class PPTTemplateParser:
    """PPT模板解析器"""

    def __init__(self, template_path):
        """
        初始化解析器

        Args:
            template_path: PPT模板文件路径
        """
        self.template_path = template_path
        self.prs = Presentation(template_path)
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height

        # 提取各页
        self.header_slide = self.prs.slides[0]  # 第一页：头图
        self.content_template = self.prs.slides[1]  # 第二页：内容页模板
        self.footer_slide = self.prs.slides[-1]  # 最后一页：尾图

        # 分析模板结构
        self.template_structure = self._analyze_template()

    def _analyze_template(self):
        """
        分析内容页模板的结构，提取款号、零售价、图片位置等信息

        Returns:
            dict: 模板结构信息
        """
        structure = {
            'title_shape': None,  # 左上角标题文本框
            'table_shape': None,  # 表格形状
            'table_info': []      # 每个单元格的信息
        }

        for shape in self.content_template.shapes:
            # 识别左上角标题文本框（非表格内的文本框）
            if shape.shape_type == 17 and not shape.has_table:  # TEXT_BOX
                if structure['title_shape'] is None:
                    structure['title_shape'] = {
                        'name': shape.name,
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                        'text': shape.text_frame.text if shape.has_text_frame else ''
                    }

            # 识别表格
            if shape.has_table:
                structure['table_shape'] = {
                    'name': shape.name,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height,
                    'rows': len(shape.table.rows),
                    'cols': len(shape.table.columns)
                }

                # 分析每个单元格
                table = shape.table
                for i, row in enumerate(table.rows):
                    for j, cell in enumerate(row.cells):
                        cell_text = cell.text.strip()
                        cell_info = {
                            'row': i,
                            'col': j,
                            'text': cell_text,
                            'is_product_code_label': TEMPLATE_LABELS['product_code'] in cell_text,
                            'is_price_label': TEMPLATE_LABELS['retail_price'] in cell_text,
                            'is_product_code_value': self._is_product_code_value(cell_text),
                            'is_price_value': self._is_price_value(cell_text),
                            'font_info': self._extract_font_info(cell)
                        }
                        structure['table_info'].append(cell_info)

        return structure

    def _is_product_code_value(self, text):
        """
        判断文本是否是款号值（纯数字或字母+数字）
        """
        if not text or TEMPLATE_LABELS['product_code'] in text:
            return False
        # 纯数字或字母+数字
        return text.isdigit() or (any(c.isalpha() for c in text) and any(c.isdigit() for c in text))

    def _is_price_value(self, text):
        """
        判断文本是否是价格值（包含价格符号+数字）
        支持全角和半角价格符号
        """
        if not text or TEMPLATE_LABELS['retail_price'] in text:
            return False
        # 包含价格符号（全角或半角）和数字
        has_symbol = TEMPLATE_LABELS['price_symbol'] in text or TEMPLATE_LABELS.get('price_symbol_fullwidth', '￥') in text
        has_digit = any(c.isdigit() for c in text)
        return has_symbol and has_digit

    def _extract_font_info(self, cell):
        """
        提取单元格字体信息
        """
        font_info = {}
        if cell.text_frame.paragraphs:
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                run = para.runs[0]
                # 安全获取颜色
                color = None
                try:
                    color = run.font.color.rgb
                except (AttributeError, TypeError):
                    color = None

                font_info = {
                    'name': run.font.name,
                    'size': run.font.size.pt if run.font.size else 9,
                    'bold': run.font.bold,
                    'italic': run.font.italic,
                    'color': color
                }
        return font_info

    def get_template_layout_info(self):
        """
        自动分析模板排版信息：
        - 每页产品数量（款号值的总数）
        - 每行产品数量（根据款号值所在列的分布）
        - 产品行数（根据款号值所在行的分布）

        Returns:
            dict: 排版信息
        """
        code_positions = []
        price_positions = []

        for cell_info in self.template_structure['table_info']:
            if cell_info['is_product_code_value']:
                code_positions.append((cell_info['row'], cell_info['col']))
            elif cell_info['is_price_value']:
                price_positions.append((cell_info['row'], cell_info['col']))

        # 按行列排序
        code_positions.sort(key=lambda x: (x[0], x[1]))

        # 分析行列分布
        code_rows = sorted(set([p[0] for p in code_positions]))
        code_cols = sorted(set([p[1] for p in code_positions]))

        # 计算每行产品数
        products_per_row = 0
        if code_rows:
            row_counts = []
            for r in code_rows:
                count = len([p for p in code_positions if p[0] == r])
                row_counts.append(count)
            products_per_row = max(row_counts) if row_counts else 0

        # 计算总行数
        num_rows = len(code_rows)

        # 计算每页产品总数
        products_per_page = len(code_positions)

        # 计算每个产品的列区域范围（用于图片居中定位）
        # 逻辑：每个产品占2列（款号列 + 其左边的标签列）
        # 图片中心基于这个2列区域的中心计算
        product_col_ranges = []
        for code_row, code_col in code_positions:
            # 款号值在奇数列(1,3,5,7)，区域=[code_col-1, code_col]
            # 款号值在偶数列(0,2,4,6)，区域=[code_col, code_col+1]
            if code_col % 2 == 1:
                col_start = code_col - 1
                col_end = code_col
            else:
                col_start = code_col
                col_end = code_col + 1
            product_col_ranges.append((code_row, col_start, col_end))

        return {
            'products_per_page': products_per_page,
            'products_per_row': products_per_row,
            'num_rows': num_rows,
            'code_positions': code_positions,
            'price_positions': price_positions,
            'product_col_ranges': product_col_ranges
        }


class PPTGenerator:
    """PPT生成器"""

    def __init__(self, template_parser):
        """
        初始化生成器

        Args:
            template_parser: PPTTemplateParser实例
        """
        self.parser = template_parser

    def generate(self, products, series_name, series_config, product_info, info_config, image_size):
        """
        生成产品手册PPT
        自动根据模板中的款号数量确定每页产品数和行列排版

        核心逻辑：
        1. 直接使用原始模板创建新的Presentation（保留所有格式）
        2. 复制内容页模板（第二页）来创建所有内容页
        3. 在每一页中填充产品数据
        4. 保留原始单元格的所有格式（字体、颜色、对齐等）

        Args:
            products: 产品列表，每个产品包含{'image_path', 'code', 'price'}
            series_name: 系列名称
            series_config: 系列名称样式配置 {'font', 'size', 'bold', 'italic'}
            product_info: 产品信息文本
            info_config: 产品信息样式配置 {'font', 'size', 'bold', 'italic'}
            image_size: 图片尺寸 {'width', 'height'}

        Returns:
            str: 生成的PPT文件路径
        """
        # 从模板自动获取排版信息
        layout_info = self.parser.get_template_layout_info()
        products_per_page = layout_info['products_per_page']

        print(f"模板分析结果：每页 {products_per_page} 个产品，每行 {layout_info['products_per_row']} 个，共 {layout_info['num_rows']} 行")

        # 直接使用原始模板创建新的Presentation，保留所有格式
        new_prs = Presentation(self.parser.template_path)

        # 获取空白布局（用于创建新幻灯片）
        blank_layout = new_prs.slide_layouts[6] if len(new_prs.slide_layouts) > 6 else new_prs.slide_layouts[-1]

        # 1. 第一页已经是头图，不需要修改

        # 2. 根据模板每页产品数，计算需要多少内容页
        total_products = len(products)
        num_pages = (total_products + products_per_page - 1) // products_per_page if products_per_page > 0 else 1

        # 3. 获取内容页模板（第二页）
        content_template_slide = new_prs.slides[1]

        # 生成所有内容页
        # 第一页直接使用模板页，后续页面复制模板页
        slides_to_fill = []

        for page_idx in range(num_pages):
            if page_idx == 0:
                # 第一页内容直接使用原来的第二页
                new_slide = content_template_slide
            else:
                # 后续页面需要复制模板页
                new_slide = self._duplicate_slide(new_prs, content_template_slide, blank_layout)

            slides_to_fill.append(new_slide)

        # 4. 填充每个页面的产品数据
        for page_idx, slide in enumerate(slides_to_fill):
            start_idx = page_idx * products_per_page
            end_idx = min(start_idx + products_per_page, total_products)
            page_products = products[start_idx:end_idx]
            self._fill_products(slide, page_products, series_name, series_config, product_info, info_config, image_size)

        # 5. 删除没有插入任何图片的空内容页（保留头图和尾图）
        self._remove_empty_content_slides(new_prs)

        # 6. 最后一页已经是尾图，不需要修改

        # 7. 保存
        output_path = os.path.join('output', 'product_manual.pptx')
        new_prs.save(output_path)

        return output_path

    def _remove_empty_content_slides(self, prs):
        """
        删除没有插入任何图片的内容页（保留头图和尾图）

        遍历所有幻灯片，检查每页是否有图片形状。
        头图（索引0）和尾图（最后一个）始终保留。
        内容页（索引1到倒数第二个）如果没有图片则删除。
        """
        p_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'

        # 获取presentation.xml中的sldIdLst元素
        sldIdLst = prs.part._element.find(f'.//{{{p_ns}}}sldIdLst')
        if sldIdLst is None:
            print("未找到sldIdLst，无法删除空页")
            return

        # 获取所有sldId元素
        all_sldIds = list(sldIdLst)
        if len(all_sldIds) <= 2:
            # 只有头图和尾图，没有内容页需要删除
            return

        # 检查每页是否有图片（排除头图和尾图）
        # 内容页范围：索引1 到 倒数第二个
        slides_to_remove = []
        for idx in range(1, len(all_sldIds) - 1):
            # 通过rId获取slide part，然后检查是否有图片
            r_ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
            sldId = all_sldIds[idx]
            r_id = sldId.get(r_ns)

            has_picture = False
            if r_id and r_id in prs.part.rels:
                rel = prs.part.rels[r_id]
                slide_part = rel.target_part
                # 检查slide的XML中是否有图片引用
                # 图片在XML中以<p:pic>元素表示
                pics = slide_part._element.findall(f'.//{{{p_ns}}}pic')
                if pics:
                    has_picture = True

            if not has_picture:
                slides_to_remove.append(idx)

        # 按从大到小的顺序删除，避免索引变化
        removed_count = 0
        for idx in sorted(slides_to_remove, reverse=True):
            sldId = all_sldIds[idx]
            sldIdLst.remove(sldId)
            removed_count += 1
            print(f"已删除空页（幻灯片索引 {idx}）")

        if removed_count > 0:
            print(f"共删除 {removed_count} 个空页")

    def _duplicate_slide(self, prs, source_slide, blank_layout):
        """
        复制幻灯片，保留所有格式

        使用python-pptx的底层XML操作来精确复制幻灯片，
        包括所有形状、表格、图片和格式设置。
        """
        # 添加新幻灯片
        new_slide = prs.slides.add_slide(blank_layout)

        # 复制幻灯片背景
        if source_slide.background.fill.type is not None:
            new_slide.background.fill.background()
            # 复制背景格式
            self._copy_background(source_slide, new_slide)

        # 复制所有形状（保持原有顺序）
        for shape in source_slide.shapes:
            self._copy_shape(new_slide, shape)

        return new_slide

    def _copy_background(self, source_slide, new_slide):
        """复制幻灯片背景"""
        try:
            # 获取源幻灯片的背景XML
            source_bg = source_slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')
            new_bg = new_slide._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld')

            if source_bg is not None and new_bg is not None:
                # 复制背景元素
                bg_elem = source_bg.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
                if bg_elem is not None:
                    # 移除新幻灯片的背景
                    existing_bg = new_bg.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}bg')
                    if existing_bg is not None:
                        new_bg.remove(existing_bg)
                    # 添加源幻灯片的背景
                    new_bg.insert(0, copy.deepcopy(bg_elem))
        except Exception as e:
            print(f"复制背景失败: {e}")

    def _copy_shape(self, new_slide, source_shape):
        """
        复制形状到新幻灯片

        对于表格，使用底层XML复制以保留所有格式（边框、填充等）
        对于文本框，保留所有文本格式
        对于图片，复制图片数据
        """
        try:
            if source_shape.has_table:
                # 表格：使用底层XML复制以保留所有格式
                self._copy_table_shape(new_slide, source_shape)
            elif source_shape.shape_type == 13:  # 图片
                self._copy_image_shape(new_slide, source_shape)
            elif source_shape.has_text_frame:
                self._copy_text_shape(new_slide, source_shape)
            else:
                # 其他形状：使用XML复制
                self._copy_shape_xml(new_slide, source_shape)
        except Exception as e:
            print(f"复制形状失败 {source_shape.name}: {e}")

    def _copy_table_shape(self, new_slide, source_shape):
        """
        复制表格形状，保留所有格式（边框、填充、字体等）

        使用底层XML操作来精确复制表格，确保所有格式都被保留。
        """
        try:
            # 获取源表格的XML元素
            source_element = source_shape._element

            # 深拷贝XML元素
            new_element = copy.deepcopy(source_element)

            # 更新新元素的ID（避免冲突）
            # 获取新的形状ID
            new_id = self._get_next_shape_id(new_slide)
            nvSpPr = new_element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}nvGraphicFramePr')
            if nvSpPr is not None:
                cNvPr = nvSpPr.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
                if cNvPr is not None:
                    cNvPr.set('id', str(new_id))

            # 将新元素添加到新幻灯片
            new_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')

        except Exception as e:
            print(f"复制表格失败: {e}")
            # 降级处理：使用常规方法复制
            self._copy_shape_xml(new_slide, source_shape)

    def _copy_image_shape(self, new_slide, source_shape):
        """复制图片形状"""
        try:
            image = source_shape.image
            image_bytes = BytesIO(image.blob)
            new_slide.shapes.add_picture(
                image_bytes,
                source_shape.left,
                source_shape.top,
                source_shape.width,
                source_shape.height
            )
        except Exception as e:
            print(f"复制图片失败: {e}")

    def _copy_text_shape(self, new_slide, source_shape):
        """复制文本形状"""
        try:
            new_shape = new_slide.shapes.add_textbox(
                source_shape.left,
                source_shape.top,
                source_shape.width,
                source_shape.height
            )
            new_frame = new_shape.text_frame
            source_frame = source_shape.text_frame

            # 复制文本框属性
            new_frame.word_wrap = source_frame.word_wrap

            for para_idx, para in enumerate(source_frame.paragraphs):
                if para_idx > 0:
                    new_frame.add_paragraph()
                new_para = new_frame.paragraphs[para_idx]
                new_para.text = para.text
                new_para.alignment = para.alignment
                new_para.level = para.level

                for run in para.runs:
                    new_run = new_para.add_run()
                    new_run.text = run.text
                    new_run.font.name = run.font.name
                    new_run.font.size = run.font.size
                    new_run.font.bold = run.font.bold
                    new_run.font.italic = run.font.italic
                    # 复制颜色
                    try:
                        if run.font.color.rgb:
                            new_run.font.color.rgb = run.font.color.rgb
                    except:
                        pass
        except Exception as e:
            print(f"复制文本框失败: {e}")

    def _copy_shape_xml(self, new_slide, source_shape):
        """使用XML复制形状"""
        try:
            source_element = source_shape._element
            new_element = copy.deepcopy(source_element)

            # 更新形状ID
            new_id = self._get_next_shape_id(new_slide)
            nvSpPr = new_element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr')
            if nvSpPr is not None:
                cNvPr = nvSpPr.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
                if cNvPr is not None:
                    cNvPr.set('id', str(new_id))

            new_slide.shapes._spTree.insert_element_before(new_element, 'p:extLst')
        except Exception as e:
            print(f"XML复制形状失败: {e}")

    def _get_next_shape_id(self, slide):
        """获取下一个可用的形状ID"""
        max_id = 0
        for shape in slide.shapes:
            try:
                shape_id = int(shape.shape_id)
                max_id = max(max_id, shape_id)
            except:
                pass
        return max_id + 1

    def _fill_products(self, slide, products, series_name, series_config, product_info, info_config, image_size):
        """
        在幻灯片中填充产品数据

        核心原则：
        1. 只修改单元格文本内容，不修改单元格格式
        2. 保留原始字体名称、大小、颜色、加粗、斜体等属性
        3. 只在用户明确指定了字体和字号时才覆盖
        4. 如果某个款号没有对应产品（无图片），清除该款号和对应价格文字
        5. 返回该页是否插入了至少一张图片

        Returns:
            bool: 该页是否插入了至少一张图片
        """
        # 找到表格
        table = None
        table_shape = None
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_shape = shape
                break

        if not table:
            return False

        # 找到左上角标题文本框并更新系列名称
        title_found = False
        for shape in slide.shapes:
            if shape.shape_type == 17 and not shape.has_table:  # TEXT_BOX
                if shape.text_frame:
                    para = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else None
                    if para and para.runs:
                        # 直接修改第一个run的文本，保留所有格式（颜色、字体等）
                        run = para.runs[0]
                        run.text = series_name
                        # 应用用户配置的字体和字号（覆盖原有）
                        run.font.name = series_config.get('font', run.font.name or 'ABC Whyte')
                        run.font.size = Pt(series_config.get('size', run.font.size.pt if run.font.size else 18))
                        run.font.bold = series_config.get('bold', run.font.bold if run.font.bold is not None else True)
                        run.font.italic = series_config.get('italic', run.font.italic if run.font.italic is not None else True)
                        # 删除多余的run
                        for extra_run in list(para.runs)[1:]:
                            extra_run._r.getparent().remove(extra_run._r)
                    else:
                        # 没有run，创建新的
                        if para:
                            run = para.add_run()
                        else:
                            run = shape.text_frame.paragraphs[0].add_run()
                        run.text = series_name
                        run.font.name = series_config.get('font', 'ABC Whyte')
                        run.font.size = Pt(series_config.get('size', 18))
                        run.font.bold = series_config.get('bold', True)
                        run.font.italic = series_config.get('italic', True)

                    title_found = True
                    break

        # 如果没有找到标题文本框，创建一个
        if not title_found and series_name:
            # 在左上角添加标题
            left = Inches(0.5)
            top = Inches(0.3)
            width = Inches(4)
            height = Inches(0.5)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = series_name
            run.font.name = series_config.get('font', 'ABC Whyte')
            run.font.size = Pt(series_config.get('size', 18))
            run.font.bold = series_config.get('bold', True)
            run.font.italic = series_config.get('italic', True)

        # 获取模板结构中的款号值位置和价格值位置
        code_value_positions = []
        price_value_positions = []

        for cell_info in self.parser.template_structure['table_info']:
            if cell_info['is_product_code_value']:
                code_value_positions.append((cell_info['row'], cell_info['col']))
            elif cell_info['is_price_value']:
                price_value_positions.append((cell_info['row'], cell_info['col']))

        # 按列排序，确保从左到右
        code_value_positions.sort(key=lambda x: (x[0], x[1]))
        price_value_positions.sort(key=lambda x: (x[0], x[1]))

        # 获取每个产品的列区域范围（用于图片居中定位）
        layout_info = self.parser.get_template_layout_info()
        product_col_ranges = layout_info.get('product_col_ranges', [])

        # 填充每个产品
        has_image_on_page = False  # 追踪该页是否有图片插入

        for idx, product in enumerate(products):
            if idx >= len(code_value_positions):
                break

            code_row, code_col = code_value_positions[idx]

            # 无论是否有图片，都填入款号和价格
            code_cell = table.cell(code_row, code_col)
            self._replace_cell_text(code_cell, product['code'])

            # 替换价格值（找到对应列的价格位置）
            for (pr, pc) in price_value_positions:
                if pc == code_col and pr > code_row:
                    price_cell = table.cell(pr, pc)
                    price_text = f"¥{product['price']}"
                    self._replace_cell_text(price_cell, price_text)
                    break

            # 检查是否有图片文件存在，有则插入
            image_path = product.get('image_path')
            if image_path and os.path.exists(image_path):
                # 插入图片（基于产品列区域范围计算居中位置）
                col_range = product_col_ranges[idx] if idx < len(product_col_ranges) else None
                self._insert_product_image(
                    slide, table_shape, table,
                    code_row, code_col,
                    image_path,
                    image_size,
                    col_range
                )
                has_image_on_page = True

        # 清除未使用的位置（超出产品数量的模板位置）
        for idx in range(len(products), len(code_value_positions)):
            code_row, code_col = code_value_positions[idx]
            self._clear_cell_and_price(table, code_row, code_col, price_value_positions)

        return has_image_on_page

    def _clear_cell_and_price(self, table, code_row, code_col, price_value_positions):
        """
        清除指定位置的款号文字和对应价格文字
        """
        code_cell = table.cell(code_row, code_col)
        self._replace_cell_text(code_cell, '')

        for (pr, pc) in price_value_positions:
            if pc == code_col and pr > code_row:
                price_cell = table.cell(pr, pc)
                self._replace_cell_text(price_cell, '')
                break

    def _copy_run_color(self, target_run, source_run):
        """
        复制source_run的颜色到target_run，支持所有颜色类型（RGB、主题色、索引色等）
        """
        try:
            source_color = source_run.font.color
            target_color = target_run.font.color

            # 获取source的颜色类型
            color_type = source_color.type

            if color_type == 1:  # RGB
                target_color.rgb = source_color.rgb
            elif color_type == 2:  # 主题色 (theme color)
                # 复制主题色元素
                from copy import deepcopy
                source_elem = source_color._color
                if source_elem is not None:
                    target_color._color = deepcopy(source_elem)
            elif color_type == 3:  # 索引色 (indexed color)
                target_color.theme_color = source_color.theme_color
            else:
                # 尝试直接复制XML元素
                from copy import deepcopy
                source_elem = source_color._color
                if source_elem is not None:
                    target_color._color = deepcopy(source_elem)
        except Exception as e:
            # 如果颜色复制失败，静默忽略（使用默认颜色）
            pass

    def _replace_cell_text(self, cell, new_text):
        """
        替换单元格文本，保留原有格式（包括字体、大小、颜色、加粗、斜体等）

        核心逻辑：
        1. 获取单元格中第一个段落的第一个run的格式
        2. 只修改run的文本内容，不清除整个段落
        3. 如果单元格有多个run，只修改第一个run的文本，清除其余run
        """
        if not cell.text_frame.paragraphs:
            return

        para = cell.text_frame.paragraphs[0]

        if para.runs:
            # 修改第一个run的文本
            para.runs[0].text = new_text
            # 删除多余的run（保留第一个run的所有格式）
            for extra_run in list(para.runs)[1:]:
                extra_run._r.getparent().remove(extra_run._r)
        else:
            # 没有现有run：新建一个
            run = para.add_run()
            run.text = new_text

    def _insert_product_image(self, slide, table_shape, table, code_row, code_col, image_path, image_size, col_range=None):
        """
        插入产品图片

        图片位置计算：
        1. 如果提供了col_range（产品列区域范围），基于该区域计算水平中心点
        2. 否则退回到款号单元格中心计算
        3. 图片下缘距离款号文字所在行上方1cm
        """
        if not os.path.exists(image_path):
            return

        try:
            # 图片尺寸（厘米转英寸）
            img_width = Inches(image_size['width'] / 2.54)
            img_height = Inches(image_size['height'] / 2.54)

            # 计算水平中心点
            if col_range:
                # 基于产品列区域范围计算中心
                _, col_start, col_end = col_range
                # 区域左边距（相对于表格）
                area_left = self._get_cell_left(table, col_start)
                # 区域右边距（相对于表格）
                area_right = self._get_cell_left(table, col_end) + self._get_cell_width(table, col_end)
                # 区域中心（相对于表格）
                area_center = area_left + (area_right - area_left) / 2
                # 转换为幻灯片绝对坐标
                img_center_x = table_shape.left + area_center
            else:
                # 降级：使用款号单元格中心
                cell_left = table_shape.left + self._get_cell_left(table, code_col)
                cell_width = self._get_cell_width(table, code_col)
                img_center_x = cell_left + cell_width / 2

            # 计算垂直位置：图片底部距离款号行顶部1cm
            code_row_top = table_shape.top + self._get_cell_top(table, code_row)
            img_bottom = code_row_top - Inches(1 / 2.54)
            img_top = img_bottom - img_height

            # 水平位置：中心对齐
            img_left = img_center_x - img_width / 2

            # 添加图片
            slide.shapes.add_picture(
                image_path,
                img_left,
                img_top,
                img_width,
                img_height
            )
        except Exception as e:
            print(f"插入图片失败 {image_path}: {e}")

    def _get_cell_left(self, table, col_idx):
        """获取指定列的左边距"""
        left = 0
        for i in range(col_idx):
            left += table.columns[i].width
        return left

    def _get_cell_top(self, table, row_idx):
        """获取指定行的上边距"""
        top = 0
        for i in range(row_idx):
            top += table.rows[i].height
        return top

    def _get_cell_width(self, table, col_idx):
        """获取指定列的宽度"""
        return table.columns[col_idx].width

    def _get_cell_height(self, table, row_idx):
        """获取指定行的高度"""
        return table.rows[row_idx].height
