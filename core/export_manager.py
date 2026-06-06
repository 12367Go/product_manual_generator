# -*- coding: utf-8 -*-
"""
导出管理模块
负责将生成的PPT转换为多种输出格式：分页图片ZIP、合并长图

核心设计原则：
1. 所有图片转换都必须基于先生成的PPT文件
2. 优先使用LibreOffice进行PPT到图片的转换（质量最高）
3. 降级方案：使用aspose.slides（如果已安装）
4. 最后方案：使用python-pptx的占位图（仅作为后备）
"""

import os
import zipfile
import shutil
import subprocess
import tempfile
from PIL import Image
from pptx import Presentation


class ExportManager:
    """导出管理器"""

    def __init__(self, output_folder):
        """
        初始化导出管理器

        Args:
            output_folder: 输出目录路径
        """
        self.output_folder = output_folder
        os.makedirs(output_folder, exist_ok=True)

    def export(self, ppt_path, formats):
        """
        导出PPT为多种格式

        Args:
            ppt_path: PPT文件路径
            formats: 导出格式列表 ['pptx', 'images_zip', 'long_image']

        Returns:
            dict: 各格式的输出文件路径
        """
        results = {}

        if 'pptx' in formats:
            results['pptx'] = ppt_path

        # 图片转换需要先确保PPT已生成
        # 无论用户是否选择pptx，都先生成PPT（这是图片转换的基础）
        if 'images_zip' in formats or 'long_image' in formats:
            # 将PPT转换为图片
            image_paths = self._convert_ppt_to_images(ppt_path)

            if 'images_zip' in formats:
                results['images_zip'] = self._create_images_zip(image_paths)

            if 'long_image' in formats:
                results['long_image'] = self._create_long_image(image_paths)

        return results

    def _convert_ppt_to_images(self, ppt_path):
        """
        将PPT转换为图片列表

        尝试多种转换方法，按优先级：
        1. LibreOffice（soffice）- 质量最高，推荐安装
        2. aspose.slides - 商业库，需要安装
        3. python-pptx占位图 - 后备方案

        Returns:
            list: 图片文件路径列表
        """
        # 创建临时目录存放图片
        temp_dir = os.path.join(self.output_folder, 'temp_images')
        os.makedirs(temp_dir, exist_ok=True)

        image_paths = []

        # 方法1：尝试使用LibreOffice转换
        try:
            image_paths = self._convert_with_libreoffice(ppt_path, temp_dir)
            if image_paths:
                print(f"LibreOffice转换成功，生成 {len(image_paths)} 张图片")
                return image_paths
        except Exception as e:
            print(f"LibreOffice转换失败: {e}")

        # 方法2：尝试使用aspose.slides
        try:
            image_paths = self._convert_with_aspose(ppt_path, temp_dir)
            if image_paths:
                print(f"Aspose转换成功，生成 {len(image_paths)} 张图片")
                return image_paths
        except Exception as e:
            print(f"Aspose转换失败: {e}")

        # 方法3：使用python-pptx生成占位图片（后备方案）
        print("使用后备方案生成图片（空白图片，仅作为占位）")
        image_paths = self._create_placeholder_images(ppt_path, temp_dir)

        return image_paths

    def _convert_with_libreoffice(self, ppt_path, output_dir):
        """
        使用LibreOffice将PPT转换为图片

        LibreOffice是开源免费的办公软件，支持将PPT导出为高质量图片。
        用户可以在 https://www.libreoffice.org/ 下载安装。

        Mac安装命令：brew install --cask libreoffice
        """
        image_paths = []

        # 检查LibreOffice是否可用
        soffice_path = self._find_soffice()
        if not soffice_path:
            print("未找到LibreOffice，请安装后重试")
            return image_paths

        # 使用LibreOffice转换
        # --headless: 无界面模式
        # --convert-to png: 转换为PNG格式
        # --outdir: 输出目录
        cmd = [
            soffice_path,
            '--headless',
            '--convert-to', 'png',
            '--outdir', output_dir,
            ppt_path
        ]

        result = subprocess.run(cmd, capture_output=True, text=True, timeout=180)

        if result.returncode == 0:
            # 收集生成的图片（按页码排序）
            base_name = os.path.splitext(os.path.basename(ppt_path))[0]
            for f in sorted(os.listdir(output_dir)):
                if f.endswith('.png') and f.startswith(base_name):
                    image_paths.append(os.path.join(output_dir, f))

        return image_paths

    def _find_soffice(self):
        """
        查找LibreOffice的可执行文件路径

        在不同操作系统上查找soffice命令的位置。
        """
        # 常见路径
        possible_paths = [
            'soffice',  # 如果在PATH中
            '/usr/bin/soffice',
            '/usr/local/bin/soffice',
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',
            '/Applications/LibreOffice.app/Contents/MacOS/soffice.bin',
            'C:\\Program Files\\LibreOffice\\program\\soffice.exe',
            'C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe',
        ]

        for path in possible_paths:
            if os.path.exists(path) or shutil.which(path):
                return path if os.path.exists(path) else shutil.which(path)

        return None

    def _convert_with_aspose(self, ppt_path, output_dir):
        """
        使用aspose.slides将PPT转换为图片

        aspose.slides是一个商业PPT处理库，支持高质量的PPT渲染。
        安装命令：pip install aspose.slides
        """
        image_paths = []

        try:
            import aspose.slides as slides

            # 加载PPT
            presentation = slides.Presentation(ppt_path)

            # 设置渲染选项
            options = slides.export.RenderingOptions()

            # 将每一页导出为图片
            for i in range(len(presentation.slides)):
                slide = presentation.slides[i]

                # 导出为图片（2倍分辨率以获得更好质量）
                bitmap = slide.get_thumbnail(2.0, 2.0)

                # 保存图片
                img_path = os.path.join(output_dir, f'page_{i+1:03d}.png')
                bitmap.save(img_path, drawing.imaging.ImageFormat.png)
                image_paths.append(img_path)

            presentation.dispose()

        except ImportError:
            print("aspose.slides未安装，跳过此方法")
        except Exception as e:
            print(f"aspose.slides转换失败: {e}")

        return image_paths

    def _create_placeholder_images(self, ppt_path, output_dir):
        """
        创建占位图片（后备方案）

        当其他转换方法都不可用时，创建与幻灯片尺寸相同的空白图片。
        这些图片不包含实际内容，仅作为占位。
        """
        prs = Presentation(ppt_path)
        image_paths = []

        # 幻灯片尺寸（转换为像素，假设96 DPI）
        width_px = int(prs.slide_width.inches * 96)
        height_px = int(prs.slide_height.inches * 96)

        for i, slide in enumerate(prs.slides):
            # 创建白色背景图片
            img = Image.new('RGB', (width_px, height_px), 'white')

            # 保存
            img_path = os.path.join(output_dir, f'page_{i+1:03d}.png')
            img.save(img_path)
            image_paths.append(img_path)

        return image_paths

    def _create_images_zip(self, image_paths):
        """
        将图片打包为ZIP文件
        """
        if not image_paths:
            return None

        try:
            zip_path = os.path.join(self.output_folder, 'product_manual_pages.zip')
            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for img_path in image_paths:
                    zipf.write(img_path, os.path.basename(img_path))
            return zip_path
        except Exception as e:
            print(f"创建ZIP失败: {e}")
            return None

    def _create_long_image(self, image_paths):
        """
        将所有图片合并为一张长图
        """
        if not image_paths:
            return None

        try:
            # 读取所有图片
            images = []
            total_height = 0
            max_width = 0

            for img_path in image_paths:
                img = Image.open(img_path)
                images.append(img)
                total_height += img.height
                max_width = max(max_width, img.width)

            # 创建长图
            long_image = Image.new('RGB', (max_width, total_height), 'white')

            # 拼接图片
            current_y = 0
            for img in images:
                # 居中放置
                x_offset = (max_width - img.width) // 2
                long_image.paste(img, (x_offset, current_y))
                current_y += img.height

            # 保存
            long_image_path = os.path.join(self.output_folder, 'product_manual_long.png')
            long_image.save(long_image_path, quality=95)

            return long_image_path

        except Exception as e:
            print(f"创建长图失败: {e}")
            return None

    def cleanup(self):
        """清理临时文件"""
        temp_dirs = [
            os.path.join(self.output_folder, 'temp_images'),
            os.path.join(self.output_folder, 'temp_long')
        ]
        for temp_dir in temp_dirs:
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir, ignore_errors=True)
